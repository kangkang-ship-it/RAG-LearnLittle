"""
Aspose.Slides Cloud API 渲染功能测试（Phase 3）

验证 AsposeCloudRenderer 是否正确调用 Aspose 模板渲染 API 并生成 PPT 文件。
每个用例输出「预想结果 vs 实际结果」对照，全部通过打印 PASS。

运行：
    python test_aspose_renderer.py

前置：
    .env 配置 ASPOSE_CLIENT_ID / ASPOSE_CLIENT_SECRET（§8.3）
    网络可达 api.aspose.cloud
"""
import io
import os
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from dotenv import load_dotenv

load_dotenv()

from pptx import Presentation
from pptx.util import Inches

from app.schemas.ppt import PPTOutline, PPTSlide
from app.services.ppt_renderer_aspose import AsposeCloudRenderer
from app.services.ppt_renderer import PythonPptxRenderer, create_renderer
from app.services.ppt_service import load_ppt_config

# ========== 测试数据（固定大纲，保证可重复对比） ==========
OUTLINE = PPTOutline(
    title="Python 装饰器全解析",
    subtitle="从语法糖到工程实践",
    slides=[
        PPTSlide(type="cover", title="Python 装饰器全解析", subtitle="从语法糖到工程实践"),
        PPTSlide(type="agenda", title="目录", items=["装饰器本质", "常见用法", "高级技巧"]),
        PPTSlide(type="section", title="装饰器本质"),
        PPTSlide(type="content", title="语法糖等价", bullets=["@decorator 等价于 func = decorator(func)", "装饰器是修改函数行为的函数", "functools.wraps 保留元信息"]),
        PPTSlide(type="summary", title="总结", bullets=["回顾完成"]),
    ],
)


def make_placeholder_template(path: str, pages: dict) -> None:
    """生成带命名页 + {{key}} 占位符的模板（Aspose 模板模式用）"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for name, texts in pages.items():
        s = prs.slides.add_slide(blank)
        s._element.cSld.set("name", name)
        for i, t in enumerate(texts):
            box = s.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(8), Inches(1))
            box.text_frame.paragraphs[0].add_run().text = t
    prs.save(path)


def collect_texts(prs: Presentation) -> str:
    return " | ".join(
        sh.text_frame.text
        for s in prs.slides
        for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    )


def run_case(no: int, name: str, expected: str, actual: str, passed: bool) -> None:
    print(f"\n【测试 {no}】{name}")
    print(f"  预想结果: {expected}")
    print(f"  实际结果: {actual}")
    print(f"  判定: {'PASS ✅' if passed else 'FAIL ❌'}")
    return passed


def main() -> int:
    print("=" * 60)
    print("Aspose.Slides Cloud API 渲染功能测试")
    print("=" * 60)
    if not os.getenv("ASPOSE_CLIENT_ID") or not os.getenv("ASPOSE_CLIENT_SECRET"):
        print("❌ 未配置 ASPOSE_CLIENT_ID/ASPOSE_CLIENT_SECRET（.env）")
        return 1

    renderer = AsposeCloudRenderer(load_ppt_config())
    results = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)

        # ========== 测试 1：模板上传到 Aspose 存储 ==========
        tmpl1 = tmp / "tmpl1.pptx"
        make_placeholder_template(str(tmpl1), {"cover": ["{{title}}", "{{subtitle}}", "{{date}}"]})
        try:
            storage_path = renderer._ensure_template_uploaded(str(tmpl1))
            # 二次调用应命中缓存（同一路径+内容）
            storage_path2 = renderer._ensure_template_uploaded(str(tmpl1))
            passed = storage_path == storage_path2 and storage_path.startswith("ppt-templates/")
            results.append(run_case(1, "模板上传到 Aspose 存储（含 md5 缓存）",
                                    "上传成功且返回存储路径；重复调用命中缓存返回同一路径",
                                    f"存储路径={storage_path}, 二次调用同一路径={storage_path == storage_path2}",
                                    passed))
        except Exception as e:
            results.append(run_case(1, "模板上传到 Aspose 存储", "上传成功",
                                    f"异常: {type(e).__name__}: {str(e)[:120]}", False))

        # ========== 测试 2：Aspose 模板渲染返回合法 PPT 文件 ==========
        try:
            data = renderer._render_with_aspose(str(tmpl1), OUTLINE, "business")
            passed = data[:2] == b"PK" and len(data) > 1000
            results.append(run_case(2, "Aspose 模板渲染返回合法 .pptx 字节流",
                                    "返回 PK 文件头、体积 > 1KB（真实 pptx 而非错误响应）",
                                    f"文件头={data[:4]!r}, 大小={len(data)}B",
                                    passed))
        except Exception as e:
            results.append(run_case(2, "Aspose 模板渲染返回合法 .pptx 字节流",
                                    "返回 PK 文件头、体积 > 1KB",
                                    f"异常: {type(e).__name__}: {str(e)[:120]}", False))
            return 0 if all(results) else 1

        # ========== 测试 3：占位符替换正确性（内容/中文/无残留） ==========
        try:
            prs = Presentation(io.BytesIO(data))
            texts = collect_texts(prs)
            passed = (
                "Python 装饰器全解析" in texts          # {{title}} 替换
                and "从语法糖到工程实践" in texts        # {{subtitle}} 替换
                and "{{" not in texts                    # 无占位符残留
                and "2026" in texts                      # {{date}} 替换（当年）
            )
            results.append(run_case(3, "占位符替换（标题/副标题/日期，中文无乱码，无残留）",
                                    "{{title}}→大纲标题、{{subtitle}}→副标题、{{date}}→日期；无 {{ 残留；中文正常",
                                    f"页数={len(prs.slides)}, 含标题={'Python 装饰器全解析' in texts}, "
                                    f"含副标题={'从语法糖到工程实践' in texts}, 含日期={'2026' in texts}, "
                                    f"无残留={'{{' not in texts}",
                                    passed))
        except Exception as e:
            results.append(run_case(3, "占位符替换正确性", "标题/副标题/日期替换且无残留",
                                    f"异常: {type(e).__name__}: {str(e)[:120]}", False))

        # ========== 测试 4：多页模板逐页填充（顶层标量键；容器循环见文档边界） ==========
        # 说明：Aspose 模板模式中独立页的占位符只匹配**顶层数据键**；
        # {{bullets}} 位于 {{sections}} 容器内（容器循环场景，需模板画 {{sections}}
        # 容器形状复制变长页），不在本用例覆盖范围（文档 §5.6 已注明）。
        tmpl4 = tmp / "tmpl4.pptx"
        make_placeholder_template(str(tmpl4), {
            "cover": ["{{title}}", "{{subtitle}}"],
            "content": ["{{date}}"],
        })
        try:
            data4 = renderer._render_with_aspose(str(tmpl4), OUTLINE, "business")
            prs4 = Presentation(io.BytesIO(data4))
            texts4 = collect_texts(prs4)
            passed = (
                len(prs4.slides) == 2                       # 模板 2 页 → 输出 2 页
                and "Python 装饰器全解析" in texts4         # 第 1 页 {{title}}
                and "从语法糖到工程实践" in texts4           # 第 1 页 {{subtitle}}
                and "2026" in texts4                        # 第 2 页 {{date}}
                and "{{" not in texts4
            )
            results.append(run_case(4, "多页模板逐页填充（顶层标量键：title/subtitle/date）",
                                    "输出页数=模板页数(2)；每页顶层占位符被大纲数据填充；无残留",
                                    f"输出页数={len(prs4.slides)}, title={'Python 装饰器全解析' in texts4}, "
                                    f"subtitle={'从语法糖到工程实践' in texts4}, date={'2026' in texts4}, "
                                    f"无残留={'{{' not in texts4}",
                                    passed))
        except Exception as e:
            results.append(run_case(4, "多页模板逐页填充", "2 页输出且内容填充",
                                    f"异常: {type(e).__name__}: {str(e)[:120]}", False))

        # ========== 测试 5：无占位符模板 → 本地 T2/T3 兜底（PPT_ENGINE=aspose_cloud 下自动降级） ==========
        tmpl5 = tmp / "tmpl5.pptx"
        make_placeholder_template(str(tmpl5), {"cover": ["静态标题"], "content": ["静态内容"]})
        # 去掉占位符（模拟纯设计模板）
        prs5 = Presentation(str(tmpl5))
        for s in prs5.slides:
            s._element.cSld.set("name", "")
        prs5.save(str(tmpl5))
        try:
            data5 = renderer.render(OUTLINE, theme="business", template_path=str(tmpl5))
            prs5 = Presentation(io.BytesIO(data5))
            texts5 = collect_texts(prs5)
            passed = (
                len(prs5.slides) == len(OUTLINE.slides)      # 大纲页数完整
                and "Python 装饰器全解析" in texts5           # 内容填充
                and "静态标题" not in texts5                   # 模板旧文本被覆盖
            )
            results.append(run_case(5, "无占位符模板 → 本地 T2/T3 兜底",
                                    "引擎自动降级本地渲染；输出页数=大纲页数；内容填充；模板旧文本无残留",
                                    f"输出页数={len(prs5.slides)}(大纲{len(OUTLINE.slides)}), "
                                    f"内容填充={'Python 装饰器全解析' in texts5}, 旧文本清除={'静态标题' not in texts5}",
                                    passed))
        except Exception as e:
            results.append(run_case(5, "无占位符模板本地兜底", "降级成功且内容正确",
                                    f"异常: {type(e).__name__}: {str(e)[:120]}", False))

    print("\n" + "=" * 60)
    passed_count = sum(results)
    print(f"测试结论: {passed_count}/{len(results)} 项通过")
    print("=" * 60)
    return 0 if passed_count == len(results) else 1


if __name__ == "__main__":
    raise SystemExit(main())
