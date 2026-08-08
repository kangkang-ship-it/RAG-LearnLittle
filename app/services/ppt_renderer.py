"""
PPT 渲染引擎（设计方案 §5）

抽象接口 + PythonPptxRenderer 实现（v1 默认）：
- 16:9 版式（13.333 × 7.5 英寸）
- 封面 / 目录 / 章节 / 内容 / 总结 五类固定版式
- 每页写入演讲者备注（讲解场景刚需）
- 风格预设（business / academic / minimal）来自 config/ppt.yaml themes
- 内容页要点超出单页上限自动分页
- Markdown 轻量清洗（去重符号、提取行首要点），不解析富文本

用户模板（§5.6）：template_path 传入时走模板分支，v1 有限支持
（Phase 1.5 spike 后实现母版/版式复用 + {{key}} 替换）；无法打开/解析
失败 → 抛错由调用方降级默认版式（三档兜底第 2 档）。

Phase 3 可选实现：AsposeCloudRenderer（模板渲染模式完整支持）。
"""
import copy
import io
import os
import re
from datetime import datetime
from typing import List, Optional, Protocol

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.core.logger_handler import logger
from app.schemas.ppt import PPTOutline, PPTSlide

# 正文字体（打开时本地渲染，无字体问题）
FONT_NAME = "微软雅黑"
CODE_FONT_NAME = "Consolas"

# 16:9 页面尺寸（英寸）
SLIDE_W, SLIDE_H = 13.333, 7.5

# 单页内容要点上限（超出自动分页）
MAX_BULLETS_PER_CONTENT = 6
# 代码块最大行数（超出折叠）
MAX_CODE_LINES = 12


def _clean_markdown(text: str) -> str:
    """轻量清洗 Markdown 语法（§5.3）：去重符号、提取行首要点，不解析富文本"""
    t = (text or "").strip()
    if not t:
        return ""
    # 行首符号（# 标题 / * - 列表 / 数字列表）
    t = re.sub(r"^(#{1,6}\s*|\*\s*|-\s*|\d+[.、]\s*)", "", t)
    # 行内标记与链接
    t = re.sub(r"[*_`#~]", "", t)
    t = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", t)
    return t[:200]


class PPTRenderer(Protocol):
    """渲染引擎抽象接口（§5.2，Phase 3 AsposeCloudRenderer 同签名）"""

    def render(
        self, outline: PPTOutline, theme: str,
        template_path: Optional[str] = None,
    ) -> bytes: ...


def create_renderer(config: Optional[dict] = None) -> PPTRenderer:
    """渲染引擎工厂（§8.3）：PPT_ENGINE=aspose_cloud 时用 Aspose，默认 python_pptx。

    AsposeCloudRenderer 内部对无占位符模板/失败场景自动降级本地，
    业务代码（PptService）无感知。
    """
    engine = os.getenv("PPT_ENGINE", "python_pptx")
    if engine == "aspose_cloud":
        from app.services.ppt_renderer_aspose import AsposeCloudRenderer

        return AsposeCloudRenderer(config)
    return PythonPptxRenderer(config)


class PythonPptxRenderer:
    """本地 python-pptx 渲染（v1 默认引擎）"""

    def __init__(self, config: Optional[dict] = None):
        self.config = config or {}
        self._themes = self.config.get("themes", {})

    # ========== 入口 ==========

    def render(
        self, outline: PPTOutline, theme: str = "business",
        template_path: Optional[str] = None,
    ) -> bytes:
        """
        渲染 PPT 大纲为 .pptx 字节流（纯同步，调用方必须放入线程池，§5.5）

        Args:
            outline: 结构化大纲（§4.4）
            theme: 风格预设（business / academic / minimal）
            template_path: 用户模板路径（§5.6）；None 或渲染失败 → 默认版式

        Returns:
            .pptx 文件字节流
        """
        # 用户模板（§5.6）：模板无法打开/解析失败 → 降级默认版式（三档兜底第 2 档）
        if template_path:
            try:
                return self._render_with_template(template_path, outline, theme)
            except NotImplementedError:
                # Phase 1 模板能力未实现（Phase 1.5 spike 后补充）→ 降级
                logger.warning("用户模板渲染能力未就绪（Phase 1.5），降级默认版式")
            except Exception as e:
                logger.warning(f"用户模板渲染失败，降级默认版式: {e}")
        return self._render_default(outline, theme)

    def _render_with_template(
        self, template_path: str, outline: PPTOutline, theme: str
    ) -> bytes:
        """模板渲染四级降级链（§5.6 扩展）：

        T1 命名页+占位符（精确模式）：模板含 cover/agenda/... 命名页 → 复制+{{key}} 填充
        T2 内容覆盖（识别+替换）：无命名页但含可识别文本 → 保留每页布局，
           识别标题/正文框并覆盖为新内容（用户「按我的模板生成」的常规路径）
        T3 母版优先（设计语言继承）：无文本或识别失败 → 保留模板母版/版式，
           用模板版式新建标准讲解页（背景/配色/字体/母版 logo 仍生效）
        → 模板文件损坏（Presentation 打不开）由 render() 外层兜底默认版式；
        → 各层内部失败逐级降级（T2 失败 → T3；T3 失败 → 默认版式）
        """
        src = Presentation(template_path)

        # ---- T1: 命名页匹配（现有逻辑，最精确） ----
        patterns: dict = {}
        for slide in src.slides:
            name = (slide._element.cSld.get("name") or "").strip().lower()
            if name in ("cover", "agenda", "section", "content", "summary"):
                patterns[name] = slide
        if patterns:
            return self._render_named_templates(src, patterns, outline, theme)

        # ---- T2: 内容覆盖（模板含可识别文本时） ----
        if self._has_any_text(src):
            try:
                return self._render_by_overlay(src, outline, theme)
            except Exception as e:
                logger.warning(f"T2 内容覆盖失败，降级 T3 母版优先: {e}")

        # ---- T3: 母版优先（设计语言继承） ----
        try:
            return self._render_master_only(src, outline, theme)
        except Exception as e:
            logger.warning(f"T3 母版优先失败，降级默认版式: {e}")
            raise

    # ========== 页重建（T1/T2 共用：原地修改模板，主题/母版/版式/背景全保留） ==========

    @staticmethod
    def _sld_id_for(src, slide):
        """通过 rId 找到 slide 对应的 sldId 元素（sldIdLst 与 slides 每次迭代包装不同）"""
        for sldId in src.slides._sldIdLst:
            rid = sldId.get(qn("r:id"))
            if rid:
                try:
                    if src.part.related_part(rid) is slide.part:
                        return sldId
                except Exception:
                    pass
        return None

    def _rebuild_slides(self, src, plan, outline: PPTOutline, theme: str) -> bytes:
        """按 plan 原地重建模板幻灯片（v1.6 修复：主题色失效）。

        plan: List[(PPTSlide, Optional[src_slide])]
            - (outline 页, 使用的模板页)：已原地覆盖内容，保留该页
            - (outline 页, None)：该页用标准讲解页（新建，基于模板版式）
        流程：新建缺失页 → 重排 sldIdLst 按 outline 顺序 → 保存。
        未在 plan 中的模板页（多余页）自动丢弃。
        """
        layout = self._pick_content_layout(src)
        colors = self._theme_colors(theme)
        sldIdLst = src.slides._sldIdLst

        ordered = []
        for s, page in plan:
            if page is None:
                slide = src.slides.add_slide(layout)
                self._render_slide(src, layout, s, colors, outline)
                ordered.append(sldIdLst[-1])
            else:
                sld_id = self._sld_id_for(src, page)
                if sld_id is not None:
                    ordered.append(sld_id)
                else:
                    # 异常兜底：模板页丢失 → 标准页
                    slide = src.slides.add_slide(layout)
                    self._render_slide(src, layout, s, colors, outline)
                    ordered.append(sldIdLst[-1])

        # 重排 sldIdLst 为 outline 顺序；
        # 未使用的模板页：drop_rel 移除其 part 引用（否则残留 part 打包时文件名冲突）
        keep = {id(el) for el in ordered}
        for el in list(sldIdLst):
            rid = el.get(qn("r:id"))
            if id(el) not in keep and rid:
                try:
                    src.part.drop_rel(rid)
                except Exception:
                    pass
            sldIdLst.remove(el)
        for el in ordered:
            sldIdLst.append(el)

        buffer = io.BytesIO()
        src.save(buffer)
        return buffer.getvalue()

    # ========== T1: 命名页 + 占位符（精确模式） ==========

    @staticmethod
    def _clone_slide_in_place(prs, src_slide):
        """同演示文稿内复制幻灯片（保留主题/版式——跨文件复制会丢主题色）。

        用于命名页被多个 outline 页复用（如模板仅 1 个 content 页、大纲 3 个 content 页）。
        """
        new_slide = prs.slides.add_slide(src_slide.slide_layout)
        for shape in list(new_slide.shapes):
            shape._element.getparent().remove(shape._element)
        for shape in src_slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape._element))
        if src_slide.has_notes_slide:
            notes_text = src_slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                new_slide.notes_slide.notes_text_frame.text = notes_text
        return new_slide

    def _render_named_templates(
        self, src, patterns: dict, outline: PPTOutline, theme: str
    ) -> bytes:
        """T1：模板幻灯片按**名称**匹配版式类型，命名页原地填充 {{key}} 占位符；
        命名页被多个大纲页复用时，同文件复制（保留主题色）"""
        # 用 list 而非 set：python-pptx Slide 对象不可哈希（定义了 __eq__）
        # set 成员测试会抛 "unhashable type: 'Slide'"（v1.6 修复）
        used: list = []
        plan = []
        for s in self._expanded_slides(outline):
            pattern = patterns.get(s.type)
            if pattern is None:
                # 该页类型模板缺失 → 标准讲解页
                logger.debug(f"模板缺失 {s.type} 版式页，该页用标准版式")
                plan.append((s, None))
                continue
            if pattern in used:
                pattern = self._clone_slide_in_place(src, pattern)
            else:
                used.append(pattern)
            # 原地填充（模板主题/母版/背景保留）
            self._fill_placeholders(pattern, self._outline_slide_mapping(s))
            if s.notes:
                pattern.notes_slide.notes_text_frame.text = s.notes
            plan.append((s, pattern))
        return self._rebuild_slides(src, plan, outline, theme)

    # ========== T2: 内容覆盖（识别 + 替换） ==========

    @staticmethod
    def _has_any_text(src) -> bool:
        """模板是否有可识别文本（无 → T2 不适用，直接 T3）"""
        for slide in src.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return True
        return False

    @staticmethod
    def _shape_text(shape) -> str:
        """形状完整文本（多段落 \n 连接）"""
        if not shape.has_text_frame:
            return ""
        return "\n".join(p.text for p in shape.text_frame.paragraphs)

    @staticmethod
    def _max_font_size(shape) -> float:
        """形状内最大 run/段落字号（pt）；无显式字号返回 0"""
        if not shape.has_text_frame:
            return 0.0
        sizes = []
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    sizes.append(r.font.size.pt)
            if p.font.size is not None:
                sizes.append(p.font.size.pt)
        return max(sizes) if sizes else 0.0

    @staticmethod
    def _is_title_placeholder(shape) -> bool:
        """版式角色：标题占位符（placeholder idx==0）——最准的标题信号（T2 规则①）"""
        try:
            return bool(shape.is_placeholder) and shape.placeholder_format.idx == 0
        except Exception:
            return False

    @staticmethod
    def _shape_top_in(shape, slide_h: float) -> float:
        """形状顶部 Y 坐标（英寸）；读取失败返回页高（视为底部）"""
        try:
            return shape.top / 914400.0
        except Exception:
            return slide_h

    def _identify_title_frame(self, slide, slide_h: float):
        """识别主标题框（T2 定稿规则）：
        ① 角色优先：标题占位符（idx==0）→ 直接命中
        ② 无占位符 → 视觉得分 = 字号 + 位置(上30%加权) + Z序(靠前加权)
           排除：文本>60字符 / 纯数字编号(如"01") / 字号<14pt / 底部10%区域
        """
        for shape in slide.shapes:
            if self._is_title_placeholder(shape):
                return shape
        candidates = []
        for idx, shape in enumerate(slide.shapes):
            text = self._shape_text(shape).strip()
            if not text:
                continue
            # spike 修复：排除纯数字/编号文本（章节大编号"01"常比真标题字号大）
            if re.fullmatch(r"\d{1,3}", text):
                continue
            size = self._max_font_size(shape)
            if len(text) > 60 or (size and size < 14):
                continue
            top = self._shape_top_in(shape, slide_h)
            if top > slide_h * 0.9:
                continue
            score = (size if size else 18) \
                + (12 if top < slide_h * 0.3 else 0) \
                + (5 if idx < 10 else 0)
            candidates.append((score, shape))
        if not candidates:
            return None
        candidates.sort(key=lambda x: -x[0])
        return candidates[0][1]

    def _identify_subtitle_frame(self, slide, title_frame, slide_h: float):
        """识别副标题框：主标题下方 0.5 英寸内、字号最大的框（T2 定稿）"""
        if title_frame is None:
            return None
        try:
            title_bottom = (title_frame.top + title_frame.height) / 914400.0
        except Exception:
            return None
        best, best_score = None, -1.0
        for shape in slide.shapes:
            if shape._element is title_frame._element or not shape.has_text_frame:
                continue
            text = self._shape_text(shape).strip()
            if not text or len(text) > 60:
                continue
            top = self._shape_top_in(shape, slide_h)
            if abs(top - title_bottom) > 0.5:
                continue
            size = self._max_font_size(shape)
            if size > best_score:
                best, best_score = shape, size
        return best

    def _identify_body_frames(self, slide, title_frame, slide_h: float) -> list:
        """识别正文框：剩余有文本框 − 页脚(底部10%) − 短文本(1-5字,非上半区) − 图注(<10pt)。

        spike 修复：用 _element 引用比较排除标题框（slide.shapes 每次迭代
        返回不同包装实例，`shape is title_frame` 恒为 False）。
        """
        title_el = title_frame._element if title_frame is not None else None
        bodies = []
        for shape in slide.shapes:
            if shape._element is title_el or not shape.has_text_frame:
                continue
            text = self._shape_text(shape).strip()
            if not text:
                continue
            top = self._shape_top_in(shape, slide_h)
            if top > slide_h * 0.9:
                continue  # 页脚/页码区
            if 1 <= len(text) <= 5 and top >= slide_h * 0.3:
                continue  # 短文本且不在上半区（页码/角标/装饰）
            size = self._max_font_size(shape)
            if size and size < 10:
                continue  # 图注/小字说明
            bodies.append(shape)
        return bodies

    @staticmethod
    def _para_is_bullet(p) -> bool:
        """段落是否有项目符号标记（buChar/buAutoNum）"""
        pPr = p._p.find(qn("a:pPr"))
        if pPr is None:
            return False
        return (pPr.find(qn("a:buChar")) is not None
                or pPr.find(qn("a:buAutoNum")) is not None)

    def _split_bullets(self, shape) -> List[str]:
        """把形状文本拆成 bullets（T2 定稿）：
        段落有项目符号或段落短(<60字) → 拆成多条；长段落 → 保留为一条
        """
        tf = shape.text_frame
        parts = []
        for p in tf.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            if self._para_is_bullet(p) or len(text) < 60:
                parts.append(text)
            else:
                parts.append(text)  # 长段落保留结构（不强行拆分）
        return parts

    def _overlay_texts(self, shape, values: List[str]) -> None:
        """覆盖形状文本（T2）：保留首个段落格式；
        单值直接替换并清空多余段落；多值按段落展开"""
        tf = shape.text_frame
        paras = list(tf.paragraphs)
        if not paras:
            return
        first_p = paras[0]
        if len(values) == 1:
            self._set_para_text(first_p._p, values[0])
            for p in paras[1:]:
                p._p.getparent().remove(p._p)
        elif values:
            for p in paras[1:]:
                p._p.getparent().remove(p._p)
            self._expand_paragraph(first_p, values)

    def _overlay_outline_slide(
        self, slide, s: PPTSlide, title_frame, slide_h: float
    ) -> bool:
        """按大纲页类型覆盖识别出的框（T2）：
        标题 → title；副标题 → subtitle（cover/section）；正文 → bullets/items。
        spike 调优：
        - content/summary 正文框 ≥ 4 个（卡片/复杂布局）→ 返回 False，
          由调用方降级 T3 标准页（多卡片无法可靠填充，不破坏布局）
        - agenda 目录卡片：每框填一个 item（保留卡片布局）
        - 其余正文框清空（避免残留模板示例文本）
        """
        # 标题
        self._overlay_texts(title_frame, [_clean_markdown(s.title)])
        # 副标题（cover / section）
        if s.subtitle:
            sub = self._identify_subtitle_frame(slide, title_frame, slide_h)
            if sub is not None:
                self._overlay_texts(sub, [_clean_markdown(s.subtitle)])
        # 正文
        bodies = self._identify_body_frames(slide, title_frame, slide_h)
        if not bodies:
            return True

        if s.type == "agenda":
            # 目录卡片：每框填一个 item（框多而 items 少时多余框清空）
            items = list(s.items or [])
            for i, body in enumerate(bodies):
                if i < len(items):
                    self._overlay_texts(body, [_clean_markdown(items[i])])
                else:
                    self._overlay_texts(body, [""])
            return True

        if s.type in ("content", "summary", "section"):
            if len(bodies) >= 4:
                # 复杂卡片布局无法可靠填充（含页序错配：section 对上了内容卡页）
                # → 该页降级 T3 标准页
                logger.debug(f"页含 {len(bodies)} 个正文框（复杂布局），降级标准版式")
                return False
            values = [_clean_markdown(b) for b in (s.bullets or [])]
            if values:
                self._overlay_texts(bodies[0], values)
            for extra in bodies[1:]:
                self._overlay_texts(extra, [""])
        # section / cover：仅标题（正文保留设计，不填）
        return True

    def _render_by_overlay(
        self, src, outline: PPTOutline, theme: str
    ) -> bytes:
        """T2 内容覆盖（识别+替换）：在模板页上**原地**覆盖文本，
        主题/母版/背景完整保留（v1.6 修复：复制到新演示文稿会丢失主题色）。

        页映射：outline 首页 → 模板第 1 页；outline 末页 → 模板末页；
        中间按序对应；模板页不足 / 无法识别标题 / 复杂布局 → 该页标准讲解页。
        """
        slide_h = src.slide_height / 914400.0
        src_slides = list(src.slides)
        outline_slides = list(self._expanded_slides(outline))
        n_src, n_out = len(src_slides), len(outline_slides)

        # 已使用的模板页索引（页映射冲突保护：模板页数 < 大纲页数时，
        # 末页映射与中间页可能撞同一模板页 → 冲突页用标准版式，避免 sldId 重复丢页）
        used_idx: set = set()
        plan = []
        for i, s in enumerate(outline_slides):
            # 页映射（v1 顺序对应）
            if i == 0:
                src_idx = 0
            elif i == n_out - 1 and n_src > 1 and n_out > 1:
                src_idx = n_src - 1
            else:
                src_idx = i
            if src_idx >= n_src or src_idx in used_idx:
                # 模板页不足 / 页映射冲突 → 该页标准讲解页
                if src_idx >= n_src:
                    logger.debug(f"模板页不足（大纲第 {i + 1} 页），该页用标准版式")
                else:
                    logger.debug(f"模板第 {src_idx + 1} 页已被占用（页映射冲突），该页用标准版式")
                plan.append((s, None))
                continue
            used_idx.add(src_idx)
            page = src_slides[src_idx]
            title_frame = self._identify_title_frame(page, slide_h)
            if title_frame is None:
                # 无法识别标题 → 该页标准讲解页（不猜错）
                logger.debug(f"模板第 {src_idx + 1} 页无法识别标题，该页用标准版式")
                plan.append((s, None))
                continue
            if not self._overlay_outline_slide(page, s, title_frame, slide_h):
                # 复杂布局无法可靠填充 → 该页标准讲解页
                plan.append((s, None))
                continue
            if s.notes:
                page.notes_slide.notes_text_frame.text = s.notes
            plan.append((s, page))

        return self._rebuild_slides(src, plan, outline, theme)

    # ========== T3: 母版优先（设计语言继承） ==========

    @staticmethod
    def _pick_content_layout(prs):
        """选含标题+正文占位符的版式（Title+Content 类），fallback 最后一个版式"""
        for layout in prs.slide_layouts:
            idxs = [ph.placeholder_format.idx for ph in layout.placeholders]
            if 0 in idxs and 1 in idxs:
                return layout
        return prs.slide_layouts[-1] if prs.slide_layouts else None

    @staticmethod
    def _drop_last_slide(prs) -> None:
        """删除演示文稿最后一页（含 slide part，避免保存时文件重名）"""
        sldIdLst = prs.slides._sldIdLst
        if not len(sldIdLst):
            return
        sldId = sldIdLst[-1]
        rId = sldId.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    def _render_master_only(
        self, src, outline: PPTOutline, theme: str
    ) -> bytes:
        """T3 母版优先：删除模板内容页（母版/版式保留），
        用模板版式新建标准讲解页——背景/配色/字体/母版 logo 全部来自模板"""
        # 删除模板全部幻灯片（含 part，避免保存时 Duplicate name；masters/layouts 保留）
        while len(src.slides._sldIdLst):
            self._drop_last_slide(src)
        layout = self._pick_content_layout(src)
        colors = self._theme_colors(theme)
        for s in self._expanded_slides(outline):
            self._render_slide(src, layout, s, colors, outline)

        buffer = io.BytesIO()
        src.save(buffer)
        return buffer.getvalue()

    # ========== 模板渲染工具 ==========

    @staticmethod
    def _outline_slide_mapping(s: PPTSlide) -> dict:
        """大纲页 → 占位符映射（标量/列表分开，供 _fill_placeholders 使用）"""
        mapping: dict = {"title": _clean_markdown(s.title)}
        if s.subtitle:
            mapping["subtitle"] = _clean_markdown(s.subtitle)
        mapping["date"] = datetime.now().strftime("%Y-%m-%d")
        if s.bullets:
            mapping["bullets"] = [_clean_markdown(b) for b in s.bullets]
        if s.items:
            mapping["items"] = [_clean_markdown(i) for i in s.items]
        if s.code:
            mapping["code"] = s.code
        return mapping

    def _fill_placeholders(self, slide, mapping: dict) -> None:
        """{{key}} 占位符替换（§5.6）：
        - 段落文本恰为 {{bullets}}/{{items}}/{{code}} → 列表展开（按项复制段落，继承格式）
        - 其余含 {{标量}} 的文本 → 行内替换（保留首个 run 格式，清空多余 run）
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in list(shape.text_frame.paragraphs):
                full = "".join(r.text for r in para.runs)
                if "{{" not in full:
                    continue
                stripped = full.strip()

                # ① 列表/代码占位符：整段展开（bullets/items 为 list，code 为 str 按行展开）
                list_key = None
                for key in ("bullets", "items", "code"):
                    if stripped == "{{" + key + "}}":
                        list_key = key
                        break
                if list_key:
                    values = mapping.get(list_key)
                    if isinstance(values, list) and values:
                        self._expand_paragraph(para, values)
                    elif isinstance(values, str) and values:
                        self._expand_paragraph(para, values.splitlines())
                    continue

                # ② 标量行内替换（如 "标题：{{title}}"；注意传 para._p 而非 para）
                for key, value in mapping.items():
                    if isinstance(value, (list, dict)):
                        continue
                    ph = "{{" + key + "}}"
                    if ph in full:
                        self._set_para_text(para._p, full.replace(ph, str(value)))
                        break

    @staticmethod
    def _expand_paragraph(para, values: List[str]) -> None:
        """把占位符段落展开为多个段落（每项一段，继承模板段落/首个 run 格式）"""
        p_el = para._p
        parent = p_el.getparent()
        idx = parent.index(p_el)
        PythonPptxRenderer._set_para_text(p_el, values[0])
        for i, value in enumerate(values[1:], start=1):
            new_p = copy.deepcopy(p_el)
            PythonPptxRenderer._set_para_text(new_p, value)
            parent.insert(idx + i, new_p)

    @staticmethod
    def _set_para_text(p_el, text: str) -> None:
        """设置段落文本：保留首个 run 的格式（rPr），删除多余 run"""
        runs = p_el.findall(qn("a:r"))
        if not runs:
            return
        t = runs[0].find(qn("a:t"))
        if t is None:
            t = runs[0].makeelement(qn("a:t"), {})
            runs[0].append(t)
        t.text = text
        for r in runs[1:]:
            p_el.remove(r)

    # ========== 默认版式 ==========

    def _render_default(self, outline: PPTOutline, theme: str) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        colors = self._theme_colors(theme)
        blank = prs.slide_layouts[6]  # blank 版式

        for slide in self._expanded_slides(outline):
            self._render_slide(prs, blank, slide, colors, outline)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _expanded_slides(self, outline: PPTOutline) -> List[PPTSlide]:
        """内容页要点超出单页上限 → 自动分页（§5.3）"""
        for s in outline.slides:
            if s.type == "content" and s.bullets and len(s.bullets) > MAX_BULLETS_PER_CONTENT:
                for i in range(0, len(s.bullets), MAX_BULLETS_PER_CONTENT):
                    yield PPTSlide(
                        type="content", title=s.title,
                        bullets=s.bullets[i:i + MAX_BULLETS_PER_CONTENT],
                        code=s.code if i == 0 else None,
                        notes=s.notes,
                    )
            else:
                yield s

    def _theme_colors(self, theme: str) -> dict:
        """取主题色；未知主题回退 business"""
        theme = theme if theme in self._themes else "business"
        t = self._themes.get(theme, {})
        return {
            "primary": t.get("primary", "1F4E79"),
            "secondary": t.get("secondary", "2E75B6"),
            "background": t.get("background", "FFFFFF"),
            "title_size": t.get("title_font_size", 36),
            "body_size": t.get("body_font_size", 20),
        }

    # ========== 版式绘制 ==========

    def _render_slide(self, prs, blank, s: PPTSlide, colors: dict, outline: PPTOutline):
        slide = prs.slides.add_slide(blank)
        # 清除版式自带的占位符形状（否则「单击此处添加标题」等提示文字
        # 会与 add_textbox 的正文图层叠加，形成重影——v1.6 修复）
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                shape._element.getparent().remove(shape._element)

        if s.type == "cover":
            # 上半部主色带 + 大标题 + 副标题 + 日期
            self._add_band(slide, 0, 0, SLIDE_W, 3.2, colors["primary"])
            self._add_textbox(slide, 0.8, 0.9, 11.7, 1.6, s.title,
                              size=44, bold=True, color="FFFFFF")
            subtitle = s.subtitle or outline.subtitle
            if subtitle:
                self._add_textbox(slide, 0.8, 2.5, 11.7, 0.6, subtitle,
                                  size=20, color="FFFFFF")
            self._add_textbox(slide, 0.8, 6.6, 11.7, 0.5,
                              datetime.now().strftime("%Y-%m-%d"),
                              size=14, color="9E9E9E")

        elif s.type == "agenda":
            self._add_textbox(slide, 0.8, 0.5, 11.7, 1.0, s.title,
                              size=32, bold=True, color=colors["primary"])
            items = s.items or [t.title for t in outline.slides if t.type == "section"]
            self._add_bullets(slide, 1.2, 1.8, 10.9, 5.0, items, size=22)

        elif s.type == "section":
            # 主色横带 + 大标题
            self._add_band(slide, 0, 2.6, SLIDE_W, 2.2, colors["primary"])
            self._add_textbox(slide, 0.8, 3.0, 11.7, 1.4, s.title,
                              size=36, bold=True, color="FFFFFF",
                              anchor=MSO_ANCHOR.MIDDLE)
            if s.subtitle:
                self._add_textbox(slide, 0.8, 4.5, 11.7, 0.6, s.subtitle,
                                  size=18, color="FFFFFF")

        elif s.type == "content":
            # 顶部辅色标题栏 + 要点 / 代码块
            self._add_band(slide, 0, 0, SLIDE_W, 0.9, colors["secondary"])
            self._add_textbox(slide, 0.6, 0.15, 12.1, 0.6, s.title,
                              size=24, bold=True, color="FFFFFF",
                              anchor=MSO_ANCHOR.MIDDLE)
            if s.code:
                self._add_code_block(slide, 0.9, 1.3, 11.5, 5.3, s.code)
            else:
                self._add_bullets(slide, 0.9, 1.4, 11.5, 5.3,
                                  s.bullets or [], size=colors["body_size"])

        elif s.type == "summary":
            self._add_textbox(slide, 0.8, 0.5, 11.7, 1.0, s.title,
                              size=32, bold=True, color=colors["primary"])
            self._add_bullets(slide, 1.2, 1.8, 10.9, 5.0,
                              s.bullets or [], size=22)

        # 每页写入演讲者备注
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    # ========== 绘制工具 ==========

    @staticmethod
    def _add_band(slide, x, y, w, h, color: str):
        """主色横带（无边框矩形）"""
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor.from_string(color)
        band.line.fill.background()

    @staticmethod
    def _add_textbox(slide, x, y, w, h, text, size=20, bold=False,
                     color="212121", align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = _clean_markdown(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT_NAME
        run.font.color.rgb = RGBColor.from_string(color)

    @staticmethod
    def _add_bullets(slide, x, y, w, h, items, size=20, color="212121"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = _clean_markdown(item)
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.name = FONT_NAME
                run.font.color.rgb = RGBColor.from_string(color)

    @staticmethod
    def _add_code_block(slide, x, y, w, h, code):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code.splitlines()[:MAX_CODE_LINES]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)
            run.font.name = CODE_FONT_NAME
            run.font.color.rgb = RGBColor.from_string("37474F")
