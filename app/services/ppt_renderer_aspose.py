"""
Aspose.Slides Cloud 渲染引擎（设计方案 §5.2 / §5.6 / §8.3，Phase 3）

模板渲染模式（高保真）：模板含 {{key}} 占位符时走 Aspose 云端渲染——
母版/版式/主题/背景完整保留，支持 {{sections}} 循环容器（变长页数）。

策略：
- 模板含 {{key}} 占位符 → Aspose 模板渲染（create_presentation_from_template）
- 无占位符模板（纯设计模板，如用户 4.pptx）→ 桥接方案 B：
  本地识别 → 注入语义键占位符 → Aspose 渲染
- Aspose 调用失败（凭证缺失/网络/配额/超时）→ 降级本地，不阻断生成

凭证：环境变量 ASPOSE_CLIENT_ID / ASPOSE_CLIENT_SECRET（§8.3，.env 配置）。

超时提示（实测校准）：Aspose 云端渲染约 40~80s（服务器处理），
桥接全链路（识别+上传+渲染+下载）约 80~100s——Plan 路径 step_timeout 已调至
180s（agent.yaml）；ReAct 路径依赖 LLM_STREAM_TIMEOUT（.env 建议 ≥120s）。
前端进度文案提示「约需 1~2 分钟」。

模板上传缓存：**模块级**按内容 md5 去重（实例级缓存对桥接临时文件
（每次 uuid 路径不同）永不命中——实测缓存失效导致重复上传/渲染）。

实施要点（真实验证校准）：
- data 必须是 XML（JSON → 服务器 500）
- name 必须是纯文件名（带文件夹 → 404）
- download_file 返回本地临时文件路径（读取后清理）
"""
import hashlib
import os
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

from app.core.logger_handler import logger
from app.schemas.ppt import PPTOutline, PPTSlide

# 模板上传缓存（模块级，进程内跨请求共享）：内容 md5 -> Aspose 存储路径
# 桥接生成的临时模板路径每次不同（uuid），实例级缓存永不命中（实测校准）
_TEMPLATE_UPLOAD_CACHE: dict = {}


class AsposeCloudRenderer:
    """Aspose.Slides Cloud 模板渲染（与 PPTRenderer 同签名，§5.2）"""

    def __init__(self, config: Optional[dict] = None):
        self.config = config or {}
        self._api = None

    # ========== 入口 ==========

    def render(
        self, outline: PPTOutline, theme: str = "business",
        template_path: Optional[str] = None,
    ) -> bytes:
        """
        渲染大纲为 .pptx 字节流（纯同步，调用方放入线程池，§5.5）

        - 无模板 → 本地默认版式
        - 模板含 {{key}} 占位符 → Aspose 模板渲染；失败 → 降级本地
        - 模板无占位符（纯设计模板）→ **桥接方案 B**：本地识别标题/正文框
          → 注入 {{slideN_*}} 唯一键占位符 → Aspose 渲染（母版背景完整继承）
          → 桥接不适用（页数不足/可注入率 < 60%）→ 本地 T2/T3
        """
        from app.services.ppt_renderer import PythonPptxRenderer

        local = PythonPptxRenderer(self.config)
        if not template_path:
            return local.render(outline, theme)

        if self._template_has_placeholders(template_path):
            try:
                return self._render_with_aspose(template_path, outline, theme)
            except Exception as e:
                logger.warning(f"Aspose 模板渲染失败，降级本地模板渲染: {e}")
            return local.render(outline, theme, template_path)

        # 桥接方案 B（无占位符纯设计模板 → Aspose）
        bridge = self._prepare_bridge_template(template_path, outline)
        if bridge is not None:
            storage_tpl, xml_data, inject_rate = bridge
            logger.info(f"引擎 aspose_cloud 桥接模式: 注入率 {inject_rate:.0%} → Aspose 渲染")
            try:
                return self._render_with_aspose_data(storage_tpl, xml_data)
            except Exception as e:
                logger.warning(f"Aspose 桥接渲染失败（注入率 {inject_rate:.0%}），降级本地: {e}")
        logger.info("引擎 aspose_cloud 桥接不适用（注入率不足/页数不匹配）→ 本地 T2/T3 渲染")
        return local.render(outline, theme, template_path)

    # ========== Aspose 模板渲染 ==========

    def _get_api(self):
        """懒初始化 SlidesApi（凭证缺失/错误 → 抛错，由 render() 降级）"""
        if self._api is None:
            from asposeslidescloud.apis.slides_api import SlidesApi
            from asposeslidescloud.configuration import Configuration

            client_id = os.getenv("ASPOSE_CLIENT_ID", "")
            client_secret = os.getenv("ASPOSE_CLIENT_SECRET", "")
            if not client_id or not client_secret:
                raise ValueError("ASPOSE_CLIENT_ID/ASPOSE_CLIENT_SECRET 未配置（设计方案 §8.3）")
            conf = Configuration()
            conf.app_sid = client_id
            conf.app_key = client_secret
            self._api = SlidesApi(conf)
        return self._api

    @staticmethod
    def _template_has_placeholders(template_path: str) -> bool:
        """模板是否含 {{key}} 占位符文本（有 → 可走 Aspose 模板模式）"""
        from pptx import Presentation

        try:
            prs = Presentation(template_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and "{{" in shape.text_frame.text:
                        return True
        except Exception:
            pass
        return False

    def _ensure_template_uploaded(self, template_path: str) -> str:
        """模板上传到 Aspose 存储（模块级内容 md5 缓存，跨请求复用），返回存储路径"""
        content = Path(template_path).read_bytes()
        digest = hashlib.md5(content).hexdigest()
        cached = _TEMPLATE_UPLOAD_CACHE.get(digest)
        if cached:
            return cached
        storage_path = f"ppt-templates/{digest}.pptx"
        with open(template_path, "rb") as f:
            self._get_api().upload_file(storage_path, f)
        _TEMPLATE_UPLOAD_CACHE[digest] = storage_path
        logger.info(f"Aspose 模板已上传: {storage_path} ({len(content)}B)")
        return storage_path

    @staticmethod
    def _xml_escape(text: str) -> str:
        return (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    @staticmethod
    def _alpha_seq(n: int) -> str:
        """0→a, 1→b, ..., 25→z, 26→aa（纯字母页面序号，避开 Aspose 数字键数组 bug）"""
        s = ""
        while True:
            n, r = divmod(n, 26)
            s = chr(97 + r) + s
            if n == 0:
                return s

    @classmethod
    def _slide_keys(cls, s: PPTSlide, section_no: int, content_no: int) -> dict:
        """大纲页 → 占位符键名前缀（纯字母，§5.6 文档）。

        实测校准：Aspose 模板模式的**数组数据不支持含数字的键名**
        （{{bullets_1}}/{{slide3_bullets}} 均不填充；纯字母 {{bullets}} 正常；
        数字键标量 {{slide1_title}} 正常）。故页面级唯一键用语义前缀 + 字母序号：
        cover / agenda / summary 固定；section_a / content_a / section_b ... 按出现序。
        """
        if s.type == "cover":
            prefix = "cover"
        elif s.type == "agenda":
            prefix = "agenda"
        elif s.type == "summary":
            prefix = "summary"
        elif s.type == "section":
            prefix = f"section_{cls._alpha_seq(section_no)}"
        else:
            prefix = f"content_{cls._alpha_seq(content_no)}"
        return {
            "type": f"{prefix}_type",
            "title": f"{prefix}_title",
            "subtitle": f"{prefix}_subtitle",
            "items": f"{prefix}_items",
            "bullets": f"{prefix}_bullets",
        }

    @classmethod
    def _outline_to_xml(cls, outline: PPTOutline) -> str:
        """outline → Aspose 模板数据（XML，模板 {{key}} 占位符契约，§5.6 文档）：

        - 顶层：{{title}} / {{subtitle}} / {{date}}（封面标量）
        - 每页语义键（页面级独立填充，键名纯字母）：
          {{cover_title}} / {{agenda_items}} / {{section_a_title}} /
          {{content_a_bullets}} / {{summary_bullets}} ...
        - 容器循环（需模板画 {{sections}} 容器形状）：{{agenda}} / {{sections}}
        列表数组用 <item> 元素（Aspose XML 数据约定）。
        """
        esc = cls._xml_escape
        parts = ['<?xml version="1.0"?>', "<Data>"]
        parts.append(f"<title>{esc(outline.title)}</title>")
        parts.append(f"<subtitle>{esc(outline.subtitle or '')}</subtitle>")
        parts.append(f"<date>{datetime.now().strftime('%Y-%m-%d')}</date>")

        section_no = content_no = 0
        for s in outline.slides:
            if s.type == "section":
                keys = cls._slide_keys(s, section_no, content_no)
                section_no += 1
            elif s.type == "content":
                keys = cls._slide_keys(s, section_no, content_no)
                content_no += 1
            else:
                keys = cls._slide_keys(s, section_no, content_no)
            parts.append(f"<{keys['type']}>{esc(s.type)}</{keys['type']}>")
            parts.append(f"<{keys['title']}>{esc(s.title)}</{keys['title']}>")
            if s.subtitle:
                parts.append(f"<{keys['subtitle']}>{esc(s.subtitle)}</{keys['subtitle']}>")
            # 实测校准：Aspose 模板模式的**数组数据（<key><item>…）不填充普通文本框**
            # （仅容器形状支持数组循环）；普通文本框填**标量含换行**（\n 分隔显示多行）。
            # 故 items/bullets 以 \n 连接的标量下发。
            if s.items:
                parts.append(
                    f"<{keys['items']}>{esc(chr(10).join(s.items))}</{keys['items']}>")
            if s.bullets:
                parts.append(
                    f"<{keys['bullets']}>{esc(chr(10).join(s.bullets))}</{keys['bullets']}>")

        # 容器循环数据（模板含 {{agenda}}/{{sections}} 容器形状时使用）
        agenda = [s.title for s in outline.slides if s.type == "section"]
        if agenda:
            parts.append("<agenda>")
            parts.extend(f"<item>{esc(t)}</item>" for t in agenda)
            parts.append("</agenda>")
        sections = [s for s in outline.slides if s.type in ("section", "content")]
        if sections:
            parts.append("<sections>")
            for s in sections:
                parts.append("<item>")
                parts.append(f"<title>{esc(s.title)}</title>")
                parts.append(f"<subtitle>{esc(s.subtitle or '')}</subtitle>")
                if s.bullets:
                    parts.append("<bullets>")
                    parts.extend(f"<item>{esc(b)}</item>" for b in s.bullets)
                    parts.append("</bullets>")
                if s.code:
                    parts.append(f"<code>{esc(s.code)}</code>")
                parts.append("</item>")
            parts.append("</sections>")

        parts.append("</Data>")
        return "\n".join(parts)

    def _render_with_aspose(
        self, template_path: str, outline: PPTOutline, theme: str
    ) -> bytes:
        """Aspose 模板渲染模式：上传模板 → XML 数据 → 生成 → 下载 → 清理。

        实施要点（真实验证校准）：data 必须 XML（JSON → 服务器 500）；
        name 纯文件名（带文件夹 → 404）；download_file 返回本地临时路径。
        """
        api = self._get_api()
        storage_tpl = self._ensure_template_uploaded(template_path)
        return self._render_with_aspose_data(storage_tpl, self._outline_to_xml(outline))

    def _render_with_aspose_data(self, storage_tpl: str, xml_data: str) -> bytes:
        """Aspose 模板模式核心：XML 数据 + 存储模板路径 → 生成 → 下载 → 清理"""
        api = self._get_api()
        out_name = f"out-{uuid.uuid4().hex}.pptx"

        api.create_presentation_from_template(
            name=out_name, template_path=storage_tpl, data=xml_data)
        try:
            result = api.download_file(out_name)
            if isinstance(result, str) and Path(result).exists():
                content = Path(result).read_bytes()
                Path(result).unlink(missing_ok=True)
            elif isinstance(result, bytes):
                content = result
            else:
                content = str(result).encode("latin-1")
            logger.info(f"Aspose 模板渲染完成: {len(content)}B")
            return content
        finally:
            try:
                api.delete_file(out_name)
            except Exception:
                pass  # 临时输出清理失败无碍

    # ========== 桥接方案 B：本地识别 → 注入占位符 → Aspose 渲染 ==========

    def _prepare_bridge_template(
        self, template_path: str, outline: PPTOutline
    ):
        """无占位符纯设计模板的 Aspose 桥接准备（方案 B）。

        流程：
        1. 打开模板副本，按 T2 页映射（首页↔页1、末页↔末页、中间按序）确定使用页
        2. 每页用本地 T2 识别（标题/副标题/正文框）→ 文本替换为唯一键占位符
           {{slideN_title}} / {{slideN_subtitle}} / {{slideN_items|bullets}}
        3. 删除未映射页（模板页 > 大纲页时）→ 页数对齐
        4. 上传模板副本 → 返回 (存储路径, XML 数据, 注入率)

        桥接判定（不满足 → 返回 None，走本地 T2/T3）：
        - 模板页数 ≥ 大纲展开页数（否则 Aspose 输出缺页）
        - 可注入页 / 大纲页数 ≥ 60%（识别失败或复杂卡片布局页过多时桥接收益低）

        返回:
            (存储路径, XML 数据, 注入率) 或 None
        """
        from pptx import Presentation

        from app.services.ppt_renderer import PythonPptxRenderer

        local = PythonPptxRenderer(self.config)
        src = Presentation(template_path)
        slide_h = src.slide_height / 914400.0
        src_slides = list(src.slides)
        outline_slides = list(outline.slides)  # 不展开分页（Aspose 页数对齐用大纲原始页数）
        n_src, n_out = len(src_slides), len(outline_slides)
        if n_src < n_out:
            logger.debug(f"桥接不适用：模板页数({n_src}) < 大纲页数({n_out})")
            return None

        # 页映射（与本地 T2 一致；冲突保护）
        used_idx: set = set()
        inject_count = 0        # 完整注入页数（标题+正文都可填，正文残留风险低）
        injectable = True
        section_no = content_no = 0
        for i, s in enumerate(outline_slides):
            if i == 0:
                src_idx = 0
            elif i == n_out - 1 and n_src > 1 and n_out > 1:
                src_idx = n_src - 1
            else:
                src_idx = i
            if src_idx >= n_src or src_idx in used_idx:
                injectable = False
                break
            used_idx.add(src_idx)
            page = src_slides[src_idx]

            if s.type == "section":
                keys = self._slide_keys(s, section_no, content_no)
                section_no += 1
            elif s.type == "content":
                keys = self._slide_keys(s, section_no, content_no)
                content_no += 1
            else:
                keys = self._slide_keys(s, section_no, content_no)

            title = local._identify_title_frame(page, slide_h)
            if title is None:
                continue  # 该页无法识别标题 → 不注入
            # 标题框 → {{xxx_title}}（_set_para_text 接收段落元素 para._p）
            title_p = title.text_frame.paragraphs[0]._p
            local._set_para_text(title_p, f"{{{{{keys['title']}}}}}")
            # 正文框（agenda → items；content/summary → bullets）
            bodies = local._identify_body_frames(page, title, slide_h)
            if len(bodies) >= 4:
                # 复杂卡片/多目录框布局：仅标题可注入（正文保留模板原样——
                # 设计模板的卡片示例内容可视为设计元素；正文残留风险记录日志）
                # 仍计入注入率（放宽判定，让卡片型模板也能走 Aspose 渲染）
                logger.debug(f"桥接：模板第 {src_idx + 1} 页复杂布局（{len(bodies)} 正文框），仅注入标题")
                inject_count += 1
                continue
            # 副标题（cover/section）
            sub = None
            if s.subtitle:
                sub = local._identify_subtitle_frame(page, title, slide_h)
                if sub is not None:
                    sub_p = sub.text_frame.paragraphs[0]._p
                    local._set_para_text(sub_p, f"{{{{{keys['subtitle']}}}}}")
            # bodies 排除副标题框（避免 {{xxx_bullets}} 覆盖已注入的 {{xxx_subtitle}}）
            if sub is not None:
                bodies = [b for b in bodies if b._element is not sub._element]
            if bodies:
                key = "items" if s.type == "agenda" else "bullets"
                body_p = bodies[0].text_frame.paragraphs[0]._p
                local._set_para_text(body_p, f"{{{{{keys[key]}}}}}")
                # 多余正文框清空（与本地 T2 一致，避免残留模板旧文本）
                for extra in bodies[1:]:
                    extra_p = extra.text_frame.paragraphs[0]._p
                    local._set_para_text(extra_p, "")
            inject_count += 1  # 完整注入页（标题+正文）

        if not injectable:
            logger.debug("桥接不适用：页映射冲突")
            return None
        inject_rate = inject_count / n_out
        if inject_rate < 0.6:
            logger.debug(f"桥接不适用：完整注入页 {inject_count}/{n_out}（{inject_rate:.0%} < 60%）")
            return None

        # 删除未映射页（模板页 > 大纲页）→ 页数对齐
        keep_idx = {self._mapped_idx(i, n_src, n_out) for i in range(n_out)}
        sldIdLst = src.slides._sldIdLst
        for pos, sldId in enumerate(list(sldIdLst)):
            if pos not in keep_idx:
                rid = sldId.get("r:id")
                if rid:
                    try:
                        src.part.drop_rel(rid)
                    except Exception:
                        pass
                sldIdLst.remove(sldId)
        # 重排为大纲顺序
        ordered = [sldIdLst[0]] if sldIdLst else []
        # 简化：删除后剩余页按原序即大纲顺序（映射是保序的），无需重排

        import tempfile as _tempfile

        tmp_path = os.path.join(_tempfile.gettempdir(), f"bridge-{uuid.uuid4().hex}.pptx")
        src.save(tmp_path)
        storage_tpl = self._ensure_template_uploaded(tmp_path)
        Path(tmp_path).unlink(missing_ok=True)
        logger.info(f"桥接准备完成: 注入 {inject_count}/{n_out} 页（{inject_rate:.0%}），模板页 {n_src}→{n_out}")
        return storage_tpl, self._outline_to_xml(outline), inject_rate

    @staticmethod
    def _mapped_idx(i: int, n_src: int, n_out: int) -> int:
        """T2 页映射（与 _render_by_overlay 一致）"""
        if i == 0:
            return 0
        if i == n_out - 1 and n_src > 1 and n_out > 1:
            return n_src - 1
        return i
